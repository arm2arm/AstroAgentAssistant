"""Helper functions for PowerPoint decks with python-pptx.
Consistent styling: backgrounds, cards, textboxes, bars, bullets.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree


# ============================================================================
# COLOR PALETTES
# ============================================================================

# Dark theme (original)
COLOR_DARK = {
    'bg': RGBColor(0x0F, 0x17, 0x2A),
    'card': RGBColor(0x1E, 0x29, 0x3B),
    'white': RGBColor(0xF1, 0xF5, 0xF9),
    'gray': RGBColor(0x94, 0xA3, 0xB8),
    'teal': RGBColor(0x14, 0xB8, 0xA6),
    'blue': RGBColor(0x3B, 0x82, 0xF6),
    'amber': RGBColor(0xF5, 0x9E, 0x0B),
    'pink': RGBColor(0xEC, 0x48, 0x99),
}

# White background + ocean blue accents
COLOR_OCEAN = {
    'bg': RGBColor(0xFF, 0xFF, 0xFF),
    'text_dark': RGBColor(0x1A, 0x2A, 0x3A),
    'text_gray': RGBColor(0x4A, 0x5A, 0x6A),
    'card_bg': RGBColor(0xF0, 0xF9, 0xFF),
    'card_border': RGBColor(0xB3, 0xE5, 0xFC),
    'deep': RGBColor(0x00, 0x3B, 0x5C),      # #003B5C
    'light': RGBColor(0x00, 0x6B, 0x8A),     # #006B8A
    'accent': RGBColor(0x00, 0xA5, 0xCF),    # #00A5CF
    'teal': RGBColor(0x00, 0x80, 0x80),      # #008080
    'cyan': RGBColor(0x48, 0xCA, 0xE4),      # #48CAE4
    'special': RGBColor(0x00, 0x2D, 0x4A),   # Very dark blue
}


# ============================================================================
# LAYOUT HELPERS
# ============================================================================

def init_deck(width=Inches(13.333), height=Inches(7.5)):
    """Create new presentation with 16:9 aspect ratio."""
    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height
    return prs


def add_blank_slide(prs):
    """Add blank slide (layout 6 = blank)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    """Set slide background color."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                               slide.shapes._spTree.part.presentation.slide_width,
                               slide.shapes._spTree.part.presentation.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False


def card(slide, x, y, w, h, fill_color=COLOR_OCEAN['card_bg'],
         border_color=COLOR_OCEAN['card_border'], border_width=Pt(1.5),
         corner_radius=0.12):
    """Create rounded rectangle card with border."""
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = fill_color
    r.line.color.rgb = border_color
    r.line.width = border_width
    r.shadow.inherit = False
    if hasattr(r.adjustments, '__setitem__'):
        r.adjustments[0] = corner_radius
    return r


def bar(slide, x, y, w, h, color):
    """Add color bar (accent line)."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


# ============================================================================
# TEXT HELPERS
# ============================================================================

def textbox(slide, x, y, w, h, text, font_size=Pt(18), color=COLOR_OCEAN['text_dark'],
            bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add textbox with optional alignment."""
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return t


def heading(slide, x, y, w, h, text, font_size=Pt(36), color=COLOR_OCEAN['special'],
            font_name="Calibri"):
    """Add bold heading."""
    return textbox(slide, x, y, w, h, text, font_size, color, bold=True,
                   font_name=font_name)


def title_big(slide, x, y, w, h, text, font_size=Pt(48),
              color=COLOR_OCEAN['special'], font_name="Calibri"):
    """Add big title (main slide title)."""
    return textbox(slide, x, y, w, h, text, font_size, color, bold=True,
                   font_name=font_name, align=PP_ALIGN.CENTER)


# ============================================================================
# LIST HELPERS
# ============================================================================

def bullet_list(slide, x, y, w, h, items, font_size=Pt(16),
                text_color=COLOR_OCEAN['text_dark'],
                bullet_char="\u2022",
                bullet_color=COLOR_OCEAN['accent'],
                space_after=Pt(8), font_name="Calibri"):
    """Add bulleted list with custom bullet color."""
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = True

    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = font_size
        p.font.color.rgb = text_color
        p.font.name = font_name
        p.space_after = space_after

        pPr = p._p.get_or_add_pPr()

        # Set bullet character
        buChar = pPr.makeelement(f"{{{ns}}}buChar", {"char": bullet_char})
        pPr.append(buChar)

        # Set bullet color
        buClr = pPr.makeelement(f"{{{ns}}}buClr", {})
        srgb = buClr.makeelement(f"{{{ns}}}srgbClr", {
            "val": f"{bullet_color[0]:02x}{bullet_color[1]:02x}{bullet_color[2]:02x}"
        })
        buClr.append(srgb)
        pPr.append(buClr)

    return t


def numbered_list(slide, x, y, w, h, items, font_size=Pt(16),
                  text_color=COLOR_OCEAN['text_dark'],
                  number_color=COLOR_OCEAN['accent'],
                  space_after=Pt(8), font_name="Calibri"):
    """Add numbered list with custom number color."""
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = True

    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    for i, item in enumerate(items, 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.text = f"{i}. {item}"
        p.font.size = font_size
        p.font.color.rgb = text_color
        p.font.name = font_name
        p.space_after = space_after

        pPr = p._p.get_or_add_pPr()

        # Number color (post-paragraph)
        buClr = pPr.makeelement(f"{{{ns}}}buClr", {})
        srgb = buClr.makeelement(f"{{{ns}}}srgbClr", {
            "val": f"{number_color[0]:02x}{number_color[1]:02x}{number_color[2]:02x}"
        })
        buClr.append(srgb)
        pPr.append(buClr)

    return t


# ============================================================================
# NUMERATION HELPER
# ============================================================================

def slide_number_oval(slide, n, x=Inches(0.6), y=Inches(0.5),
                      oval_size=Inches(0.55), color=COLOR_OCEAN['light']):
    """Add oval circle with slide number inside."""
    # Oval shape
    cr = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, oval_size, oval_size)
    cr.fill.solid()
    cr.fill.fore_color.rgb = color
    cr.line.fill.background()
    cr.shadow.inherit = False

    # Number text (centered in oval)
    t = textbox(slide, x, y + Inches(0.08), oval_size, Oval_size * 0.8,
                str(n), Pt(20), RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                align=PP_ALIGN.CENTER)
    return cr


# ============================================================================
# TOP ACCENT BAR
# ============================================================================

def top_accent_bar(slide, color=COLOR_OCEAN['accent'], height=Inches(0.08)):
    """Add full-width accent bar at top of slide."""
    width = slide.shapes._spTree.part.presentation.slide_width
    return bar(slide, 0, 0, width, height, color)


# ============================================================================
# SAVING
# ============================================================================

def save_pdf(pptx_path, pdf_path=None):
    """Convert PPTX to PDF using LibreOffice headless mode."""
    import subprocess

    if pdf_path is None:
        pdf_path = pptx_path.replace('.pptx', '.pdf')

    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'pdf', pptx_path],
        capture_output=True, text=True, cwd='/tmp'
    )

    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

    return pdf_path


# ============================================================================
# EXAMPLE USAGE (commented out)
# ============================================================================

# def example_deck():
#     """Example: Create simple deck with ocean blue palette."""
#     prs = init_deck()
#
#     # Slide 1: Title
#     s = add_blank_slide(prs)
#     bg(s, COLOR_OCEAN['bg'])
#     top_accent_bar(s, COLOR_OCEAN['accent'])
#     title_big(s, Inches(1.5), Inches(2), Inches(10), Inches(2),
#               "My Presentation")
#     heading(s, Inches(1.5), Inches(4.5), Inches(10), Inches(1),
#             "Subtitle or Tagline", font_size=Pt(24))
#
#     # Slide 2: Problem + Vision
#     s = add_blank_slide(prs)
#     bg(s, COLOR_OCEAN['bg'])
#     heading(s, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
#             "Overview")
#
#     # Problem card
#     card(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3))
#     heading(s, Inches(1.2), Inches(1.7), Inches(4.5), Inches(0.4),
#             "Problem", font_size=Pt(20), color=COLOR_OCEAN['teal'])
#     bullet_list(s, Inches(1.2), Inches(2.2), Inches(5), Inches(2),
#                 ["Point one", "Point two", "Point three"])
#
#     # Vision card
#     card(s, Inches(6.8), Inches(1.5), Inches(5.5), Inches(3))
#     heading(s, Inches(7.2), Inches(1.7), Inches(4.5), Inches(0.4),
#             "Vision", font_size=Pt(20), color=COLOR_OCEAN['light'])
#     bullet_list(s, Inches(7.2), Inches(2.2), Inches(5), Inches(2),
#                 ["Goal one", "Goal two", "Goal three"])
#
#     prs.save("example.pptx")
#     return "example.pptx"
